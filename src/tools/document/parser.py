#!/usr/bin/env python3
# Путь: /mnt/ai_data/ai-agent/src/tools/document/parser.py
"""Парсер документов различных форматов"""

from pathlib import Path
from typing import Any, List, Union, IO, cast
import PyPDF2
from docx import Document  # type: ignore[import-untyped]
import openpyxl  # type: ignore[import-untyped]
from pptx import Presentation  # type: ignore[import-untyped]
import markdown  # type: ignore[import-untyped]
from loguru import logger


class DocumentParser:
    """Парсер для работы с документами"""
    
    def __init__(self, config: Any) -> None:
        self.config = config
        logger.info("📄 DocumentParser инициализирован")
    
    def parse(self, file_path: Union[str, Path]) -> str:
        """
        Парсинг документа в зависимости от расширения
        
        Args:
            file_path: путь к файлу (строка или Path)
            
        Returns:
            текст документа
        """
        # Преобразуем в Path для работы с методами
        path = Path(file_path)
        
        if not path.exists():
            logger.error(f"❌ Файл не найден: {file_path}")
            return ""
        
        ext = path.suffix.lower()
        
        try:
            if ext == '.pdf':
                return self._parse_pdf(path)
            elif ext in ['.docx', '.doc']:
                return self._parse_docx(path)
            elif ext in ['.xlsx', '.xls']:
                return self._parse_xlsx(path)
            elif ext in ['.pptx', '.ppt']:
                return self._parse_pptx(path)
            elif ext == '.md':
                return self._parse_md(path)
            elif ext == '.txt':
                return self._parse_txt(path)
            else:
                logger.warning(f"⚠️ Неподдерживаемый формат: {ext}")
                return ""
                
        except Exception as e:
            logger.error(f"❌ Ошибка парсинга {file_path}: {e}")
            return ""
    
    def _parse_pdf(self, file_path: Path) -> str:
        """Парсинг PDF файла"""
        try:
            text: List[str] = []
            with open(file_path, 'rb') as f:
                pdf_reader = PyPDF2.PdfReader(f)
                for page in pdf_reader.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text.append(page_text)
            return '\n'.join(text)
        except Exception as e:
            logger.error(f"❌ Ошибка парсинга PDF: {e}")
            return ""
    
    def _parse_docx(self, file_path: Path) -> str:
        """Парсинг DOCX файла"""
        try:
            # Исправлено: Document ожидает str или IO, а не Path (ошибка arg-type)
            doc = Document(str(file_path))
            return '\n'.join([paragraph.text for paragraph in doc.paragraphs if paragraph.text])
        except Exception as e:
            logger.error(f"❌ Ошибка парсинга DOCX: {e}")
            return ""
    
    def _parse_xlsx(self, file_path: Path) -> str:
        """Парсинг XLSX файла"""
        try:
            wb = openpyxl.load_workbook(str(file_path), data_only=True)
            text: List[str] = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_text = ' '.join([str(cell) for cell in row if cell])
                    if row_text:
                        text.append(row_text)
            return '\n'.join(text)
        except Exception as e:
            logger.error(f"❌ Ошибка парсинга XLSX: {e}")
            return ""
    
    def _parse_pptx(self, file_path: Path) -> str:
        """Парсинг PPTX файла"""
        try:
            prs = Presentation(str(file_path))
            text: List[str] = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text.append(shape.text)
            return '\n'.join(text)
        except Exception as e:
            logger.error(f"❌ Ошибка парсинга PPTX: {e}")
            return ""
    
    def _parse_md(self, file_path: Path) -> str:
        """Парсинг MD файла"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            html = markdown.markdown(content)
            return cast(str, html)
        except Exception as e:
            logger.error(f"❌ Ошибка парсинга MD: {e}")
            return ""
    
    def _parse_txt(self, file_path: Path) -> str:
        """Парсинг TXT файла"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
            logger.error(f"❌ Ошибка парсинга TXT: {e}")
            return ""
